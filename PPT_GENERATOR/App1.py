import os
import streamlit as st
from dotenv import load_dotenv
import google.generativeai as genai
from pptx import Presentation
import io
import re

# Load environment variables
load_dotenv()
GEMINI_API_KEY = os.getenv("GOOGLE_API_KEY")

# Configure the GenAI model
genai.configure(api_key=GEMINI_API_KEY)

def generate_content(title, topics):
    """
    Generate content using the Gemini model from genai.
    """
    if not GEMINI_API_KEY:
        st.error("Gemini API key not found. Please set it in the .env file.")
        return "Error: Missing API key."
    
    prompt = f"Generate ppt for the topic: {', '.join(topics)} under title {title}"
    
    try:
        model = genai.GenerativeModel('gemini-pro')  # Using the GenAI model
        response = model.generate_content(prompt)
        return response.text  # Returning the generated text
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return "Error generating content."

def split_into_slides(content):
    """
    Split the generated content into slides based on the 'Slide 1', 'Slide 2', etc., and fix formatting issues
    while keeping **bold** formatting intact around slide titles.
    """
    slides = []
    slide_pattern = re.compile(r"(Slide \d+:.*?)(?=Slide \d+:|$)", re.DOTALL)  # Match Slide 1, Slide 2, etc.
    matches = slide_pattern.findall(content)

    for match in matches:
        match = match[:-2] 
        slides.append(f"**{match}")
    
    return slides

def remove_double_asteris(input_list):
    # Create a new list where each string has the double asteris removed
    result = [item.replace('**', '') for item in input_list]
    return result


def create_presentation(title, topics):
    content = generate_content(title, topics)
    if content.startswith("Error"):
        return content  # Returning the error if there was one
    
    # Split the content into individual slides
    slides = split_into_slides(content)
    filtered_slides = remove_double_asteris(slides)
    st.text_area(f"raw content",content,height=200)
    return filtered_slides

def process_slide_title(slide):
    # Split the slide content by the newline to separate the slide number, title, and points
    parts = slide.split("\n", 1)  # Split into two parts: slide and the rest of the content
    
    if len(parts) == 2:
        # Extract the slide number and title from the first part
        slide_header = parts[0].strip()
        # Match slide number and title using regex
        match = re.match(r"Slide (\d+): (.+)", slide_header)
        
        if match:
            slide_number = match.group(1)  # Extract slide number
            title = match.group(2).strip()  # Extract title
            
            # Extract the content (points) after the title
            points = parts[1].strip()
            
            # Split points by newline and clean up each point, removing the bullet ('*') and extra spaces
            points_list = [point.strip("* ").strip() for point in points.split("\n") if point.strip()]
            
            # Join the points into a single string, separated by commas
            points_string = ", ".join(points_list)
            
            return slide_number, title, points_string
    
    return None

def generate_ppt(slides):
    prs = Presentation()
    for i, slide_content in enumerate(slides, 1):
        no,title,content = process_slide_title(slide_content)
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use a Title and Content layout
        title_placeholder = slide.shapes.title
        content_placeholder = slide.shapes.placeholders[1]  # Body placeholder for content
        # Set title for the slide
        title_placeholder.text = f"{no}.{title}"
        # Set the content for the slide
        content_placeholder.text = content
    # Save the presentation to a binary stream
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# Streamlit UI
st.title("AI-Powered Presentation Creator")

# Input form
with st.form("presentation_form"):
    title = st.text_input("Enter the presentation title", "Gen AI")
    topics_input = st.text_area("Enter topics (one per line)", "Introduction to GenAI\nApplications of GenAI")
    submit = st.form_submit_button("Generate Presentation")

if submit:
    topics = [topic.strip() for topic in topics_input.split("\n") if topic.strip()]
    
    if not topics:
        st.error("Please enter at least one topic.")
    else:
        st.info("Generating your presentation. Please wait...")
        slides = create_presentation(title, topics)
        
        if isinstance(slides, str) and slides.startswith("Error"):
            st.error(slides)
        else:
            st.success("Presentation content generated successfully!")
            
            # Create PPT and provide a download link
            ppt_stream = generate_ppt(slides)
            for i, slide in enumerate(slides, 1):
                st.subheader(f"Slide {i}")
                st.text_area(f"Slide {i} Content", slide, height=200)
            st.download_button(
                label="Download PowerPoint Presentation",
                data=ppt_stream,
                file_name=f"{title.replace(' ', '_')}_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )