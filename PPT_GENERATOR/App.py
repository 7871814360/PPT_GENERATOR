import os
import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv
import google.generativeai as genai

# Load environment variables
load_dotenv()
GEMINI_API_KEY = os.getenv("GOOGLE_API_KEY")

# Configure the GenAI model
genai.configure(api_key=GEMINI_API_KEY)

def generate_content(prompt):
    """
    Generate content using the Gemini model from genai.
    """
    if not GEMINI_API_KEY:
        st.error("Gemini API key not found. Please set it in the .env file.")
        return "Error: Missing API key."
    
    try:
        model = genai.GenerativeModel('gemini-pro')  # Using the GenAI model
        response = model.generate_content(prompt)
        return response.text  # Returning the generated text
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return "Error generating content."

def set_font_size(placeholder, size: int):
    """
    Set font size for all paragraphs in a placeholder.
    """
    for paragraph in placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)

def split_content(content, max_length=700):
    """
    Split long content into smaller chunks if it exceeds the max length.
    """
    chunks = []
    while len(content) > max_length:
        # Find the last space within the max_length limit to split the content
        split_index = content.rfind(" ", 0, max_length)
        if split_index == -1:
            split_index = max_length  # if no space, just split at max_length
        chunks.append(content[:split_index])
        content = content[split_index:].strip()
    chunks.append(content)  # Append the last chunk
    return chunks

def create_presentation(title, topics):
    """
    Create a presentation using pptx and Gemini API (via GenAI model).
    """
    prs = Presentation()
    
    # Add a title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    subtitle_placeholder = slide.placeholders[1]
    subtitle_placeholder.text = "Generated by AI"
    
    # Reduce font size for title and subtitle
    set_font_size(title_placeholder, 32)  # Title font size
    set_font_size(subtitle_placeholder, 20)  # Subtitle font size
    content = generate_content(f"Generate ppt for the topic:{topics} under title {title}")
    # Add slides for each topic
    for topic in topics:
        # Generate the content for the topic
        # content = generate_content(f"Write detailed information about {topic}.")
        
        
        # Display content in Streamlit (before generating PowerPoint)
        st.write(f"### {topic}")
        st.write(content)  # Display the full generated content
        
        # Split the content if it's too long
        content_chunks = split_content(content, max_length=900)
        
        # Create the first slide with the first chunk of content
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide_title = slide.shapes.title
        slide_title.text = topic
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content_chunks[0]  # First chunk of content
        
        # Remove bullets from content
        for paragraph in content_placeholder.text_frame.paragraphs:
            # Ensure there are no bullets
            paragraph.level = 0
            paragraph.bullet = None
            paragraph.alignment = PP_ALIGN.LEFT  # Optional: Align text left

        # Reduce font size for slide title and content
        set_font_size(slide_title, 24)  # Heading font size
        set_font_size(content_placeholder, 18)  # Content font size
        
        # Create additional slides for each remaining chunk of content
        for chunk in content_chunks[1:]:  # Skip the first chunk (already added)
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide_title = slide.shapes.title
            slide_title.text = topic
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = chunk  # New chunk of content
            
            # Remove bullets from content
            for paragraph in content_placeholder.text_frame.paragraphs:
                # Ensure there are no bullets
                paragraph.level = 3
                paragraph.bullet = None
                paragraph.alignment = PP_ALIGN.LEFT  # Optional: Align text left
            
            set_font_size(slide_title, 24)  # Heading font size
            set_font_size(content_placeholder, 18)  # Content font size

    # Save the presentation
    file_name = f"{title.replace(' ', '_')}.pptx"
    prs.save(file_name)

    # Return the file name
    return file_name


# Streamlit UI
st.title("AI-Powered Presentation Creator")

# Input form
with st.form("presentation_form"):
    title = st.text_input("Enter the presentation title", "Gen AI")
    topics_input = st.text_area("Enter topics (one per line)", "Introduction to GenAI\nApplications of GenAI")
    submit = st.form_submit_button("Generate Presentation")

if submit:
    topics = [topic.strip() for topic in topics_input.split("\n") if topic.strip()]
    
    if not topics:
        st.error("Please enter at least one topic.")
    else:
        st.info("Generating your presentation. Please wait...")
        file_name = create_presentation(title, topics)
        st.success(f"Presentation '{file_name}' created successfully!")
        
        # Show download button
        with open(file_name, "rb") as file:
            st.download_button(
                label="Download Presentation",
                data=file,
                file_name=file_name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )